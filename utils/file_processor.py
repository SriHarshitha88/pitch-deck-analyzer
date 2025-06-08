import os
from typing import List, Dict, Union
import PyPDF2
from pptx import Presentation
import streamlit as st

class FileProcessor:
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from a PDF file."""
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text

    @staticmethod
    def extract_text_from_ppt(file_path: str) -> str:
        """Extract text from a PowerPoint file."""
        text = ""
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    @staticmethod
    def process_uploaded_file(uploaded_file) -> Dict[str, Union[str, str]]:
        """Process an uploaded file and return its content and metadata."""
        # Create temporary file
        temp_path = f"temp_{uploaded_file.name}"
        with open(temp_path, "wb") as f:
            f.write(uploaded_file.getvalue())

        # Extract text based on file type
        file_extension = os.path.splitext(uploaded_file.name)[1].lower()
        if file_extension == '.pdf':
            content = FileProcessor.extract_text_from_pdf(temp_path)
        elif file_extension in ['.ppt', '.pptx']:
            content = FileProcessor.extract_text_from_ppt(temp_path)
        else:
            raise ValueError("Unsupported file format. Please upload PDF or PPT/PPTX files.")

        # Clean up temporary file
        os.remove(temp_path)

        return {
            "filename": uploaded_file.name,
            "content": content,
            "file_type": file_extension
        }

    @staticmethod
    def save_report(report_content: str, filename: str, format: str = "txt") -> str:
        """Save the generated report to a file."""
        reports_dir = "reports"
        os.makedirs(reports_dir, exist_ok=True)
        
        output_path = os.path.join(reports_dir, f"{filename}.{format}")
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(report_content)
        
        return output_path 